"""
generate_ppt.py — Fill the official Redrob PPT template with IntelliRank content.
Preserves all backgrounds/banners. Only adds/replaces content text boxes.

Run: python generate_ppt.py
Output: output/IntelliRank_Deck.pptx  (then export as PDF for submission)
"""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

TEMPLATE = Path("Idea Submission Template _ Redrob.pptx")
DATA_JSON = Path("frontend/public/candidates_data.json")
OUTPUT = Path("output/IntelliRank_Deck.pptx")

# ── Colours ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
BLUE   = RGBColor(0x00, 0x84, 0xFF)
LBLUE  = RGBColor(0x31, 0x9A, 0xFF)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x4A, 0x60, 0x80)
GREEN  = RGBColor(0x10, 0xB9, 0x81)

EMU_IN = 914400  # 1 inch in EMU


def inches(x): return int(x * EMU_IN)


def add_textbox(slide, left, top, width, height, text, font_size=9,
                bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True,
                line_spacing=None):
    txb = slide.shapes.add_textbox(
        inches(left), inches(top), inches(width), inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or NAVY
    if line_spacing:
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 1000)))
    return txb


def add_multiline(slide, left, top, width, height, lines, base_size=9,
                  base_color=None, wrap=True):
    """Add a text box with multiple paragraphs, each being a (text, size, bold, color) tuple."""
    txb = slide.shapes.add_textbox(
        inches(left), inches(top), inches(width), inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, bold, color = item, base_size, False, base_color or NAVY
        else:
            text, size, bold, color = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or NAVY
    return txb


def add_box(slide, left, top, width, height, fill_rgb, text="", text_color=None,
            font_size=9, bold=False, align=PP_ALIGN.CENTER):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt as Pt2
    from pptx.oxml.ns import nsmap
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        inches(left), inches(top), inches(width), inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.width = Pt(1)
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = text_color or WHITE
    return shape


def remove_questions(slide):
    """Clear existing question-text content box (keep title box and background)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            # Large content boxes contain the questions — detect by multi-line prompts
            if "?" in txt and shape.width > inches(5):
                tf = shape.text_frame
                # Keep the shape but clear text
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ""
                for i in range(len(tf.paragraphs) - 1, 0, -1):
                    p = tf.paragraphs[i]._p
                    p.getparent().remove(p)
                tf.paragraphs[0].runs[0].text = "" if tf.paragraphs[0].runs else ""


# ─────────────────────────────────────────────────────────────────────────────
def main():
    with open(DATA_JSON, encoding="utf-8") as f:
        data = json.load(f)
    top = data["candidates"][:10]

    prs = Presentation(str(TEMPLATE))
    slides = prs.slides

    # ── SLIDE 1: Title / Team Info ────────────────────────────────────────────
    s = slides[0]
    for shape in s.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text
            if "Team Name" in txt:
                tf = shape.text_frame
                tf.paragraphs[0].runs[0].text = "Team Name :  IntelliRank"
                tf.paragraphs[0].runs[0].font.color.rgb = WHITE
                tf.paragraphs[0].runs[0].font.bold = True
                tf.paragraphs[0].runs[0].font.size = Pt(14)
            elif "Team Leader" in txt:
                tf = shape.text_frame
                tf.paragraphs[0].runs[0].text = "Team Leader Name :  Pranav"
                tf.paragraphs[0].runs[0].font.color.rgb = WHITE
                tf.paragraphs[0].runs[0].font.size = Pt(12)
            elif "Problem Statement" in txt:
                tf = shape.text_frame
                tf.paragraphs[0].runs[0].text = "Problem Statement :  Data & AI Challenge — Intelligent Candidate Discovery"
                tf.paragraphs[0].runs[0].font.color.rgb = WHITE
                tf.paragraphs[0].runs[0].font.size = Pt(11)

    # ── SLIDE 2: Solution Overview ────────────────────────────────────────────
    remove_questions(slides[1])
    add_multiline(slides[1], 0.41, 1.38, 9.30, 3.80, [
        ("What is IntelliRank?", 11, True, BLUE),
        ("A multi-signal AI ranking engine that scores all 100,000 candidates across 4 independent signals "
         "and fuses them into a ranked shortlist with per-candidate explanations — in under 35 seconds on a CPU.", 9, False, NAVY),
        ("", 5, False, NAVY),
        ("4 Scoring Signals:", 10, True, BLUE),
        ("  Skill Match (35%)  —  Cluster coverage × proficiency × duration × endorsements", 9, False, NAVY),
        ("  Career Fit  (30%)  —  Experience Gaussian, title relevance, company type, location", 9, False, NAVY),
        ("  Behavioral  (15%)  —  Availability, responsiveness, reliability, engagement", 9, False, NAVY),
        ("  Semantic    (20%)  —  TF-IDF bigram cosine similarity to JD query (offline, no API)", 9, False, NAVY),
        ("", 5, False, NAVY),
        ("What differentiates IntelliRank from keyword matching?", 10, True, BLUE),
        ("  Graceful degradation — missing signals excluded & weights renormalized (never zeroed)", 9, False, NAVY),
        ("  Proficiency-weighted skills — Expert + 36 months + 25 endorsements > bare listing", 9, False, NAVY),
        ("  Behavioral signals are first-class — a fast-response, immediately available candidate ranks higher", 9, False, NAVY),
        ("  Explicit disqualifiers — consulting-only careers and wrong-domain expertise penalized directly", 9, False, NAVY),
    ])

    # ── SLIDE 3: JD Understanding & Evaluation ────────────────────────────────
    remove_questions(slides[2])
    add_multiline(slides[2], 0.41, 1.38, 4.50, 3.80, [
        ("Key Requirements Extracted from JD", 10, True, BLUE),
        ("Must-have skill clusters (5):", 9, True, NAVY),
        ("  1. Embeddings & Retrieval  (sentence-transformers, BGE, E5)", 8.5, False, NAVY),
        ("  2. Vector Databases  (FAISS, Pinecone, Weaviate, Qdrant, Milvus)", 8.5, False, NAVY),
        ("  3. Python  (strong production code)", 8.5, False, NAVY),
        ("  4. Ranking Evaluation  (NDCG, MRR, MAP, A/B testing)", 8.5, False, NAVY),
        ("  5. NLP / Information Retrieval  background", 8.5, False, NAVY),
        ("", 4, False, NAVY),
        ("Experience band:  5–9 years  (ideal = 7)", 8.5, False, NAVY),
        ("Location:  Pune · Noida · HYD · MUM · DEL · BLR  (Tier-1)", 8.5, False, NAVY),
        ("", 4, False, NAVY),
        ("Disqualifiers:", 9, True, NAVY),
        ("  Consulting-only (TCS/Infosys/Wipro etc.)", 8.5, False, GREY),
        ("  CV / speech / robotics expertise", 8.5, False, GREY),
        ("  No production code in 18+ months", 8.5, False, GREY),
    ])
    add_multiline(slides[2], 5.10, 1.38, 4.60, 3.80, [
        ("How We Evaluate Beyond Keywords", 10, True, BLUE),
        ("Skill scoring formula per cluster:", 9, True, NAVY),
        ("  score = proficiency_mult × 0.5", 8.5, False, NAVY),
        ("         + duration_mult   × 0.3", 8.5, False, NAVY),
        ("         + endorsement_mult× 0.2", 8.5, False, NAVY),
        ("", 4, False, NAVY),
        ("Proficiency:  Expert=1.0  Adv=0.85  Mid=0.65  Beg=0.40", 8, False, NAVY),
        ("Duration:     36m+=1.0   18-36=0.85  6-18=0.65  <6=0.40", 8, False, NAVY),
        ("Endorsements: 20+=1.0   10-20=0.85  3-10=0.70  <3=0.50", 8, False, NAVY),
        ("", 4, False, NAVY),
        ("Semantic engine: TF-IDF (50K bigram features)", 8.5, False, NAVY),
        ("Sees 'dense retrieval' = 'embeddings retrieval'", 8.5, False, NAVY),
        ("Works fully offline — no network calls during ranking", 8.5, False, GREEN),
    ])

    # ── SLIDE 4: Ranking Methodology ─────────────────────────────────────────
    remove_questions(slides[3])
    # Formula box
    add_box(slides[3], 0.41, 1.38, 9.30, 0.60,
            NAVY, "final_score = 0.35×skill  +  0.30×career  +  0.15×behavioral  +  0.20×semantic  (weights renormalized if signal missing)",
            font_size=9, bold=True)
    add_multiline(slides[3], 0.41, 2.08, 4.50, 3.00, [
        ("Retrieve", 10, True, BLUE),
        ("All 100K candidates loaded into memory from candidates.jsonl.", 8.5, False, NAVY),
        ("TF-IDF matrix (100K×50K sparse) precomputed once and cached.", 8.5, False, NAVY),
        ("No vector DB needed at PoC scale — NumPy cosine similarity.", 8.5, False, NAVY),
        ("", 4, False, NAVY),
        ("Score  (per candidate, vectorized)", 10, True, BLUE),
        ("  Skill Match: cluster coverage × proficiency/duration/endorsement weights", 8.5, False, NAVY),
        ("  Career Fit:  Gaussian(exp_years, mu=7, sigma=2) × title × location × company_type", 8.5, False, NAVY),
        ("  Behavioral: availability × responsiveness × reliability × engagement", 8.5, False, NAVY),
        ("  Semantic:   cosine_sim(TF-IDF(candidate), TF-IDF(JD_query)) × 100", 8.5, False, NAVY),
    ])
    add_multiline(slides[3], 5.10, 2.08, 4.60, 3.00, [
        ("Rank & Output", 10, True, BLUE),
        ("Sort by final_score DESC, candidate_id ASC (tie-break).", 8.5, False, NAVY),
        ("Top-100 exported as CSV + XLSX with reasoning column.", 8.5, False, NAVY),
        ("", 4, False, NAVY),
        ("Weight Presets (tunable via UI):", 9, True, NAVY),
        ("  Balanced:     skill=35%  career=30%  sem=20%  beh=15%", 8, False, NAVY),
        ("  Senior Hire:  skill=30%  career=40%  sem=15%  beh=15%", 8, False, NAVY),
        ("  IC / Hacker:  skill=40%  career=25%  sem=20%  beh=15%", 8, False, NAVY),
        ("", 4, False, NAVY),
        ("Runtime: 33 seconds  (100K candidates, CPU-only, no GPU)", 9, True, GREEN),
    ])

    # ── SLIDE 5: Explainability & Data Validation ─────────────────────────────
    remove_questions(slides[4])
    add_multiline(slides[4], 0.41, 1.38, 4.50, 3.80, [
        ("How Ranking Decisions Are Explained", 10, True, BLUE),
        ("Fully deterministic templated reasoning — no LLM, no hallucination risk.", 8.5, False, NAVY),
        ("", 4, False, NAVY),
        ("Example output (rank #1):", 9, True, NAVY),
        ("\"Staff Machine Learning Engineer with 7.0 yrs exp; covers", 8, False, GREY),
        ("5/5 must-have clusters (embeddings/retrieval, vector/db,", 8, False, GREY),
        ("python); India-based (Kochi); open to work; notice 60d;", 8, False, GREY),
        ("response rate 95%; final score 72.60\"", 8, False, GREY),
        ("", 4, False, NAVY),
        ("Each field in the reasoning maps directly to a raw data", 8.5, False, NAVY),
        ("field — no inference, no fabrication.", 8.5, False, NAVY),
        ("Score breakdown panel in UI shows 4 sub-scores with bars + radar chart.", 8.5, False, NAVY),
    ])
    add_multiline(slides[4], 5.10, 1.38, 4.60, 3.80, [
        ("Data Quality & Validation", 10, True, BLUE),
        ("Missing signals:", 9, True, NAVY),
        ("  Signal excluded → weights renormalized over remaining signals.", 8.5, False, NAVY),
        ("  Never scored as 0 (would unfairly penalize sparse profiles).", 8.5, False, NAVY),
        ("", 4, False, NAVY),
        ("Wrong-domain detection:", 9, True, NAVY),
        ("  Explicit keyword list: computer vision, speech, robotics.", 8.5, False, NAVY),
        ("  Each wrong-domain keyword found → -10% skill penalty.", 8.5, False, NAVY),
        ("", 4, False, NAVY),
        ("Consulting firm penalty:", 9, True, NAVY),
        ("  Pure consulting career (TCS/Infosys/etc.) → company_type_score = 0.15", 8.5, False, NAVY),
        ("  Mixed consulting+product → 0.45  |  All product → 0.95", 8.5, False, NAVY),
        ("", 4, False, NAVY),
        ("Suspicious profiles filtered by low behavioral composite (< 10).", 8.5, False, NAVY),
    ])

    # ── SLIDE 6: End-to-End Workflow ──────────────────────────────────────────
    remove_questions(slides[5])
    steps = [
        ("1", "JD Config", "5 must-have clusters,\nexp range, disqualifiers,\nlocation prefs"),
        ("2", "Load & Normalize", "100K candidates.jsonl\n→ adapter.py\n→ structured profiles"),
        ("3", "TF-IDF Precompute", "Offline once\n421MB sparse matrix\ncached to disk"),
        ("4", "Score All", "skill + career\n+ behavioral\n+ semantic (TF-IDF)"),
        ("5", "Fuse & Rank", "Weighted average\n+ renormalize\n→ sort top-100"),
        ("6", "Output", "CSV + XLSX\n+ Web UI\n+ Firestore log"),
    ]
    box_w = 1.45
    gap = 0.10
    start_x = 0.35
    for idx, (num, title, detail) in enumerate(steps):
        x = start_x + idx * (box_w + gap)
        add_box(slides[5], x, 1.45, box_w, 0.38, BLUE, f"Step {num}: {title}", font_size=8, bold=True)
        add_textbox(slides[5], x, 1.85, box_w, 1.30, detail, font_size=7.5, color=NAVY)

        if idx < len(steps) - 1:
            add_textbox(slides[5], x + box_w + 0.01, 1.55, 0.10, 0.30, "→", font_size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    add_multiline(slides[5], 0.35, 3.25, 9.30, 1.90, [
        ("Key Design Decisions:", 9, True, BLUE),
        ("  Fully offline ranking — no network calls during scoring (submission spec compliant).", 8.5, False, NAVY),
        ("  Pre-compute once, serve static JSON — Firebase Hosting delivers results in <2s (no backend needed).", 8.5, False, NAVY),
        ("  Single adapter.py — all raw field name changes isolated to one file; rest of pipeline unchanged.", 8.5, False, NAVY),
        ("  Missing-signal renormalization — never zeros a candidate unfairly.", 8.5, False, NAVY),
    ])

    # ── SLIDE 7: System Architecture ─────────────────────────────────────────
    s7 = slides[6]
    # Left column: pipeline
    add_box(s7, 0.30, 1.00, 2.60, 0.40, NAVY, "candidates.jsonl  (100K records, 465 MB)", font_size=8, bold=True)
    add_textbox(s7, 1.50, 1.42, 0.30, 0.30, "↓", 12, True, BLUE, PP_ALIGN.CENTER)
    add_box(s7, 0.30, 1.70, 2.60, 0.38, LBLUE, "adapter.py  — normalize()", font_size=8, bold=False)
    add_textbox(s7, 1.50, 2.10, 0.30, 0.30, "↓", 12, True, BLUE, PP_ALIGN.CENTER)
    add_box(s7, 0.30, 2.38, 2.60, 0.38, LBLUE, "scorer.py  — 4 sub-scores", font_size=8, bold=False)
    add_textbox(s7, 1.50, 2.78, 0.30, 0.30, "↓", 12, True, BLUE, PP_ALIGN.CENTER)
    add_box(s7, 0.30, 3.06, 2.60, 0.38, LBLUE, "fuse_scores()  — weighted average", font_size=8, bold=False)

    # Middle: TF-IDF
    add_box(s7, 3.20, 1.00, 2.80, 0.40, NAVY, "TF-IDF Cache  (precomputed, offline)", font_size=8, bold=True)
    add_textbox(s7, 4.45, 1.42, 0.30, 0.30, "↓", 12, True, BLUE, PP_ALIGN.CENTER)
    add_box(s7, 3.20, 1.70, 2.80, 0.38, LBLUE, "tfidf_semantic.py  — cosine sim", font_size=8, bold=False)
    add_textbox(s7, 4.45, 2.10, 0.30, 0.30, "↓", 12, True, BLUE, PP_ALIGN.CENTER)
    add_box(s7, 3.20, 2.38, 2.80, 0.38, LBLUE, "semantic_score  → scorer pipeline", font_size=8, bold=False)

    # Right: JD config
    add_box(s7, 6.30, 1.00, 2.80, 0.40, NAVY, "jd_config.py  (JD parsed offline)", font_size=8, bold=True)
    add_textbox(s7, 7.55, 1.42, 0.30, 0.30, "↓", 12, True, BLUE, PP_ALIGN.CENTER)
    add_box(s7, 6.30, 1.70, 2.80, 0.38, LBLUE, "Cluster weights, disqualifiers, prefs", font_size=8, bold=False)

    # Bottom: merge → output
    add_textbox(s7, 1.50, 3.46, 0.30, 0.30, "↓", 12, True, BLUE, PP_ALIGN.CENTER)
    add_box(s7, 0.30, 3.74, 9.30, 0.42, BLUE, "Sort → Top-100  →  export_candidates_json.py  →  candidates_data.json  (420 KB)", font_size=9, bold=True)
    add_textbox(s7, 1.50, 4.18, 0.30, 0.25, "↓", 11, True, BLUE, PP_ALIGN.CENTER)
    add_box(s7, 0.30, 4.43, 4.45, 0.38, RGBColor(0x10, 0xB9, 0x81), "Firebase Hosting  — React UI (no backend)", font_size=8, bold=True)
    add_box(s7, 4.85, 4.43, 4.75, 0.38, RGBColor(0x10, 0xB9, 0x81), "Firestore  — ranking run logs", font_size=8, bold=True)

    # ── SLIDE 8: Results & Performance ───────────────────────────────────────
    remove_questions(slides[7])
    # Top results table
    add_box(slides[7], 0.41, 1.38, 9.20, 0.32, NAVY, "Top 10 Ranked Candidates  (from 100,000)", font_size=9, bold=True)

    col_headers = ["Rank", "Title", "Exp", "Location", "Open", "Notice", "Score"]
    col_xs = [0.41, 0.85, 3.70, 4.30, 5.70, 6.20, 7.00]
    col_ws = [0.42, 2.80, 0.58, 1.35, 0.48, 0.75, 1.61]
    for i, (h, x, w) in enumerate(zip(col_headers, col_xs, col_ws)):
        add_box(slides[7], x, 1.72, w, 0.26, LBLUE, h, font_size=7.5, bold=True)
    for row_i, c in enumerate(top):
        y = 1.72 + 0.26 + row_i * 0.24
        vals = [
            f"#{row_i+1}",
            c.get("current_title","")[:28],
            f"{c.get('years_experience',0):.1f}y",
            (c.get("location","") or c.get("country","")).split(",")[0][:16],
            "Yes" if c.get("open_to_work") else "No",
            f"{c.get('notice_period_days',90)}d",
            f"{c.get('final_score',0):.2f}",
        ]
        fill = RGBColor(0xF0, 0xF7, 0xFF) if row_i % 2 == 0 else WHITE
        for j, (val, x, w) in enumerate(zip(vals, col_xs, col_ws)):
            box_fill = fill
            box_color = NAVY
            bold = False
            if j == 6:  # Score column
                sc = c.get("final_score", 0)
                box_fill = BLUE if sc >= 65 else LBLUE
                box_color = WHITE
                bold = True
            add_box(slides[7], x, y, w, 0.22, box_fill, val, text_color=box_color, font_size=7, bold=bold)

    # Performance stats
    add_multiline(slides[7], 0.41, 4.28, 9.20, 1.00, [
        ("Performance:", 9, True, BLUE),
        ("  Runtime: 33s for 100K candidates on CPU (Win11, 8-core, 16GB RAM, no GPU)  |  Web demo: <5s (pre-scored JSON on Firebase CDN)  |  Top score: 72.60  |  Score range: 48.81–72.60", 8, False, NAVY),
    ])

    # ── SLIDE 9: Technologies ─────────────────────────────────────────────────
    remove_questions(slides[8])
    tech_rows = [
        ("Layer", "Technology", "Why"),
        ("Scoring Pipeline", "Python 3.11 + NumPy + scikit-learn", "Fast vectorized ops, no GPU"),
        ("Semantic Engine", "TF-IDF (sklearn, 50K bigrams)", "Fully offline, no model download"),
        ("Skill Scoring", "Custom rule engine (scorer.py)", "Transparent, tunable, deterministic"),
        ("Data Normalization", "adapter.py (single file)", "All field changes isolated here"),
        ("API", "FastAPI (local dev)", "Auto OpenAPI docs, async-ready"),
        ("Frontend Framework", "React 19 + Vite + TypeScript", "Fast build, modern DX"),
        ("Styling", "Tailwind CSS + Framer Motion", "Responsive, smooth animations"),
        ("Charts", "Recharts", "React-native, SVG-based"),
        ("Hosting", "Firebase Hosting (static CDN)", "Global edge, free tier"),
        ("State Persistence", "Firestore (GCP)", "Stores ranking run history"),
        ("AI Assistance", "Claude Sonnet 4.6 (Anthropic)", "Architecture, pipeline, debugging"),
    ]
    col_w3 = [2.0, 3.8, 3.4]
    col_x3 = [0.41, 2.43, 6.25]
    hdr_colors = [NAVY, NAVY, NAVY]
    for ci, (h, x, w) in enumerate(zip(tech_rows[0], col_x3, col_w3)):
        add_box(slides[8], x, 1.38, w, 0.28, NAVY, h, font_size=9, bold=True)
    for ri, row in enumerate(tech_rows[1:]):
        y = 1.68 + ri * 0.27
        for ci, (val, x, w) in enumerate(zip(row, col_x3, col_w3)):
            fill = RGBColor(0xF0, 0xF7, 0xFF) if ri % 2 == 0 else WHITE
            add_box(slides[8], x, y, w, 0.25, fill, val, text_color=NAVY, font_size=8)

    # ── SLIDE 10: Submission Assets ───────────────────────────────────────────
    s10 = slides[9]
    for shape in s10.shapes:
        if shape.has_text_frame and "Github" in shape.text_frame.text:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            for i in range(len(tf.paragraphs)-1, 0, -1):
                tf.paragraphs[i]._p.getparent().remove(tf.paragraphs[i]._p)
            tf.paragraphs[0].runs[0].text = "" if tf.paragraphs[0].runs else ""

    add_multiline(s10, 0.41, 1.38, 9.20, 3.80, [
        ("GitHub Repository  (public)", 10, True, BLUE),
        ("  https://github.com/pranayukey200/Redrob-Data-AI-_challenge", 9, False, NAVY),
        ("", 5, False, NAVY),
        ("Live Demo  (hosted on Firebase)", 10, True, BLUE),
        ("  https://intellisense-2253d.web.app", 9, False, NAVY),
        ("  — Paste or select the JD, choose weight preset, click Find Best-Fit Candidates", 8.5, False, GREY),
        ("  — Loads 100K pre-scored candidates, re-ranks in browser in <1s", 8.5, False, GREY),
        ("", 5, False, NAVY),
        ("Submission Files", 10, True, BLUE),
        ("  output/submission.csv      — top-100 ranked candidates  (candidate_id, rank, score, reasoning)", 9, False, NAVY),
        ("  output/ranked_candidates.xlsx  — formatted XLSX with all 15 columns + colour-coded scores", 9, False, NAVY),
        ("  output/IntelliRank_Deck.pptx   — this presentation (export as PDF for submission)", 9, False, NAVY),
        ("", 5, False, NAVY),
        ("Contact:  pranayuworkspace@gmail.com  |  Team: IntelliRank  |  Solo submission", 8, False, GREY),
    ])

    OUTPUT.parent.mkdir(exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved {OUTPUT}  ({OUTPUT.stat().st_size/1024:.0f} KB)")
    print("Open it in PowerPoint/LibreOffice, export to PDF for submission.")

if __name__ == "__main__":
    main()

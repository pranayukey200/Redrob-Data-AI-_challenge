"""
Professional PPT redesign with links and improved layout
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

# Load existing presentation
prs = Presentation('output/IntelliRank_Deck_Enhanced.pptx')

# Create new presentation
new_prs = Presentation()
new_prs.slide_width = prs.slide_width
new_prs.slide_height = prs.slide_height

# Get blank layout
blank_layout = new_prs.slide_layouts[6]

def add_title_shape(slide, text, top, size=36, color=(15, 32, 56)):
    """Add a title textbox"""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.8))
    frame = box.text_frame
    frame.text = text
    para = frame.paragraphs[0]
    para.font.size = Pt(size)
    para.font.bold = True
    para.font.color.rgb = RGBColor(*color)
    para.alignment = PP_ALIGN.CENTER
    return box

def add_text_shape(slide, text, left, top, width, height, size=12, align=PP_ALIGN.LEFT, color=(0,0,0)):
    """Add a text textbox"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.text = text
    para = frame.paragraphs[0]
    para.font.size = Pt(size)
    para.font.color.rgb = RGBColor(*color)
    para.alignment = align
    return box

# ========== SLIDE 1: TITLE SLIDE ==========
slide1 = new_prs.slides.add_slide(blank_layout)
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(15, 32, 56)

# Logo/Title
title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
title_frame = title_box.text_frame
title_frame.text = 'IntelliRank'
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(66)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(255, 255, 255)
title_para.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.3), Inches(8), Inches(0.8))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = 'AI-Powered Candidate Discovery & Ranking'
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(22)
subtitle_para.font.color.rgb = RGBColor(200, 220, 255)
subtitle_para.alignment = PP_ALIGN.CENTER

# Challenge
challenge_box = slide1.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(7), Inches(0.6))
challenge_frame = challenge_box.text_frame
challenge_frame.text = 'Redrob Data & AI Challenge: Intelligent Candidate Discovery'
challenge_para = challenge_frame.paragraphs[0]
challenge_para.font.size = Pt(16)
challenge_para.font.italic = True
challenge_para.font.color.rgb = RGBColor(120, 200, 255)
challenge_para.alignment = PP_ALIGN.CENTER

# Team
team_box = slide1.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.5))
team_frame = team_box.text_frame
team_frame.text = 'Team IntelliRank | Solo Participant: Pranav'
team_para = team_frame.paragraphs[0]
team_para.font.size = Pt(14)
team_para.font.color.rgb = RGBColor(180, 180, 180)
team_para.alignment = PP_ALIGN.CENTER

# ========== SLIDE 2: QUICK ACCESS LINKS ==========
slide2 = new_prs.slides.add_slide(blank_layout)

add_title_shape(slide2, 'Quick Access', 0.4, 32)

# Links section
links_data = [
    ('🌐 Live Demo', 'https://intellisense-2253d.web.app/', 'Try the ranking system with sample data'),
    ('💻 GitHub Repository', 'https://github.com/pranayukey200/Redrob-Data-AI-_challenge', 'Full source code, documentation, and reproducibility'),
    ('📊 Sandbox Environment', 'https://intellisense-2253d.web.app/', 'Firebase-hosted interactive demo'),
]

y_pos = 1.5
for emoji_title, url, desc in links_data:
    # Title
    title_box = slide2.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = emoji_title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(20)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(15, 32, 56)

    # URL
    url_box = slide2.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.4), Inches(7.5), Inches(0.3))
    url_frame = url_box.text_frame
    url_frame.text = url
    url_para = url_frame.paragraphs[0]
    url_para.font.size = Pt(12)
    url_para.font.color.rgb = RGBColor(0, 102, 204)
    url_para.font.underline = True

    # Description
    desc_box = slide2.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.7), Inches(7.5), Inches(0.3))
    desc_frame = desc_box.text_frame
    desc_frame.text = desc
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(11)
    desc_para.font.color.rgb = RGBColor(100, 100, 100)

    y_pos += 1.4

# Key stats box
stats_box = slide2.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.5))
stats_frame = stats_box.text_frame
stats_frame.text = 'System Performance'
stats_frame.paragraphs[0].font.size = Pt(18)
stats_frame.paragraphs[0].font.bold = True
stats_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

p = stats_frame.add_paragraph()
p.text = '⚡ 69 seconds runtime  |  🎯 100K candidates ranked  |  💻 CPU-only  |  🔒 No network calls'
p.font.size = Pt(12)
p.alignment = PP_ALIGN.CENTER

# ========== SLIDE 3: PROBLEM STATEMENT ==========
slide3 = new_prs.slides.add_slide(blank_layout)

add_title_shape(slide3, 'The Challenge', 0.4, 32)

# Challenge description
challenge_box = slide3.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(1.2))
challenge_frame = challenge_box.text_frame
challenge_frame.word_wrap = True

p1 = challenge_frame.paragraphs[0]
p1.text = 'Rank 100,000 candidates for a Senior AI Engineer role at Redrob AI (Series A startup)'
p1.font.size = Pt(18)
p1.font.bold = True
p1.space_after = Pt(12)

p2 = challenge_frame.add_paragraph()
p2.text = 'Requirements: Top 100 best-fit candidates with explainable reasoning, completed in under 5 minutes on CPU-only hardware with no external API calls.'
p2.font.size = Pt(14)
p2.level = 0

# Key constraints
constraints_box = slide3.shapes.add_textbox(Inches(1), Inches(2.8), Inches(3.8), Inches(3))
constraints_frame = constraints_box.text_frame

c_title = constraints_frame.paragraphs[0]
c_title.text = 'Constraints'
c_title.font.size = Pt(20)
c_title.font.bold = True
c_title.font.color.rgb = RGBColor(15, 32, 56)
c_title.space_after = Pt(10)

constraints = [
    '⏱️  Runtime: ≤ 5 minutes',
    '💾  Memory: ≤ 16 GB RAM',
    '🚫  No GPU allowed',
    '🌐  No network during ranking',
    '📊  100 rows with reasoning',
]

for const in constraints:
    p = constraints_frame.add_paragraph()
    p.text = const
    p.font.size = Pt(13)
    p.space_after = Pt(6)

# Evaluation criteria
eval_box = slide3.shapes.add_textbox(Inches(5.2), Inches(2.8), Inches(3.8), Inches(3))
eval_frame = eval_box.text_frame

e_title = eval_frame.paragraphs[0]
e_title.text = 'Evaluation Criteria'
e_title.font.size = Pt(20)
e_title.font.bold = True
e_title.font.color.rgb = RGBColor(15, 32, 56)
e_title.space_after = Pt(10)

criteria = [
    '📈  NDCG@10, NDCG@50, MAP',
    '🎯  P@10 precision',
    '🧠  Reasoning quality',
    '⚠️  Honeypot detection',
    '🏗️  Code reproducibility',
]

for crit in criteria:
    p = eval_frame.add_paragraph()
    p.text = crit
    p.font.size = Pt(13)
    p.space_after = Pt(6)

# ========== SLIDE 4: OUR SOLUTION OVERVIEW ==========
slide4 = new_prs.slides.add_slide(blank_layout)

add_title_shape(slide4, 'Our Solution: IntelliRank', 0.4, 32)

# Solution description
sol_box = slide4.shapes.add_textbox(Inches(1), Inches(1.3), Inches(8), Inches(1))
sol_frame = sol_box.text_frame
sol_frame.word_wrap = True
sol_frame.text = 'A multi-signal AI ranking engine that scores all 100,000 candidates across 4 independent signals and fuses them into a ranked shortlist with per-candidate explanations.'
sol_para = sol_frame.paragraphs[0]
sol_para.font.size = Pt(16)
sol_para.alignment = PP_ALIGN.CENTER

# Four signals
signals_y = 2.5
signals = [
    ('Skill Match', '35%', 'Cluster coverage × proficiency × duration × endorsements'),
    ('Career Fit', '30%', 'Experience years, title relevance, location, company type'),
    ('Semantic', '20%', 'TF-IDF cosine similarity to JD (50K bigram features)'),
    ('Behavioral', '15%', 'Availability, responsiveness, reliability, engagement'),
]

colors = [
    RGBColor(74, 144, 226),   # Blue
    RGBColor(82, 196, 26),    # Green
    RGBColor(250, 173, 20),   # Yellow
    RGBColor(235, 75, 60),    # Red
]

x_positions = [0.8, 3.0, 5.2, 7.4]

for i, (name, weight, desc) in enumerate(signals):
    # Signal box
    signal_box = slide4.shapes.add_textbox(Inches(x_positions[i]), Inches(signals_y), Inches(2), Inches(2))
    signal_frame = signal_box.text_frame
    signal_frame.word_wrap = True

    # Name
    name_para = signal_frame.paragraphs[0]
    name_para.text = name
    name_para.font.size = Pt(16)
    name_para.font.bold = True
    name_para.font.color.rgb = colors[i]
    name_para.alignment = PP_ALIGN.CENTER
    name_para.space_after = Pt(4)

    # Weight
    weight_para = signal_frame.add_paragraph()
    weight_para.text = weight
    weight_para.font.size = Pt(22)
    weight_para.font.bold = True
    weight_para.font.color.rgb = colors[i]
    weight_para.alignment = PP_ALIGN.CENTER
    weight_para.space_after = Pt(8)

    # Description
    desc_para = signal_frame.add_paragraph()
    desc_para.text = desc
    desc_para.font.size = Pt(9)
    desc_para.alignment = PP_ALIGN.CENTER

# Key differentiators
diff_box = slide4.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1.5))
diff_frame = diff_box.text_frame

diff_title = diff_frame.paragraphs[0]
diff_title.text = 'What Makes IntelliRank Different'
diff_title.font.size = Pt(16)
diff_title.font.bold = True
diff_title.alignment = PP_ALIGN.CENTER
diff_title.space_after = Pt(10)

diff_text = diff_frame.add_paragraph()
diff_text.text = '✓ Graceful degradation with missing signals  |  ✓ Proficiency-weighted skills  |  ✓ Behavioral signals as first-class  |  ✓ Explicit disqualifier penalties'
diff_text.font.size = Pt(12)
diff_text.alignment = PP_ALIGN.CENTER

# Copy remaining slides from original
for i in range(3, len(prs.slides)):
    slide_layout = prs.slide_layouts[6]
    new_slide = new_prs.slides.add_slide(slide_layout)

    for shape in prs.slides[i].shapes:
        el = shape.element
        new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

# Save
new_prs.save('output/IntelliRank_Deck_Final.pptx')
print('Created professional PPT: output/IntelliRank_Deck_Final.pptx')
print('New slides:')
print('  1. Title slide with branding')
print('  2. Quick Access Links (GitHub, Demo, Sandbox)')
print('  3. Problem Statement & Challenge')
print('  4. Solution Overview with 4 signals')
print('  5-onwards: Original content (methodology, results, etc.)')
